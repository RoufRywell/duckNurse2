from django.shortcuts import render
from django.http import HttpResponse, FileResponse
from django.core.files.storage import FileSystemStorage
from docx import Document
from reportlab.pdfgen import canvas
from reportlab.lib.pagesizes import letter
from reportlab.pdfbase import pdfmetrics
from reportlab.lib.styles import getSampleStyleSheet, ParagraphStyle
from reportlab.platypus import SimpleDocTemplate, Paragraph, Spacer, PageBreak, Image, Table, TableStyle
from reportlab.lib.units import inch
from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE
from PyPDF2 import PdfReader
import os
import io
import re
from PIL import Image as PILImage
import fitz  # PyMuPDF

def fix_text_formatting(text):
    """Metni düzeltir ve kelime birleştirme sorunlarını çözer."""
    # Yaygın kelime birleştirme sorunlarını düzelt
    replacements = {
        r'\bthe\b': ' the ',
        r'\band\b': ' and ',
        r'\ban\b': ' an ',
        r'\ba\b': ' a ',
        r'\bof\b': ' of ',
        r'\bin\b': ' in ',
        r'\bon\b': ' on ',
        r'\bat\b': ' at ',
        r'\bto\b': ' to ',
        r'\bfor\b': ' for ',
        r'\bwith\b': ' with ',
        r'\bby\b': ' by ',
        r'\bfrom\b': ' from ',
        r'\bas\b': ' as ',
        r'\bis\b': ' is ',
        r'\bare\b': ' are ',
        r'\bwas\b': ' was ',
        r'\bwere\b': ' were ',
        r'\bbe\b': ' be ',
        r'\bhave\b': ' have ',
        r'\bhas\b': ' has ',
        r'\bhad\b': ' had ',
        r'\bdo\b': ' do ',
        r'\bdoes\b': ' does ',
        r'\bdid\b': ' did ',
        r'\bwill\b': ' will ',
        r'\bwould\b': ' would ',
        r'\bshall\b': ' shall ',
        r'\bshould\b': ' should ',
        r'\bcan\b': ' can ',
        r'\bcould\b': ' could ',
        r'\bmay\b': ' may ',
        r'\bmight\b': ' might ',
        r'\bmust\b': ' must ',
    }
    
    for pattern, replacement in replacements.items():
        text = re.sub(pattern, replacement, text, flags=re.IGNORECASE)
    
    # Fazla boşlukları temizle
    text = re.sub(r'\s+', ' ', text)
    return text.strip()

def extract_text_from_word(file_path):
    """Word dosyasından metni çıkarır."""
    try:
        doc = Document(file_path)
        text = []
        for paragraph in doc.paragraphs:
            if paragraph.text.strip():
                text.append(fix_text_formatting(paragraph.text.strip()))
        return "\n\n".join(text)
    except Exception as e:
        return str(e)

def extract_text_from_powerpoint(file_path):
    """PowerPoint dosyasından metni çıkarır."""
    try:
        prs = Presentation(file_path)
        text = []

        for slide in prs.slides:
            for shape in slide.shapes:
                if hasattr(shape, "text"):
                    if shape.text.strip():
                        text.append(fix_text_formatting(shape.text.strip()))

        return "\n\n".join(text)
    except Exception as e:
        return str(e)

def extract_text_from_pdf(file_path):
    """PDF dosyasından metni çıkarır."""
    try:
        reader = PdfReader(file_path)
        text = []
        for page in reader.pages:
            content = page.extract_text()
            if content.strip():
                text.append(fix_text_formatting(content.strip()))
        return "\n\n".join(text)
    except Exception as e:
        return str(e)

def extract_images_from_powerpoint(file_path):
    """PowerPoint'ten resimleri çıkarır."""
    try:
        prs = Presentation(file_path)
        images = []
        temp_dir = os.path.join(os.path.dirname(file_path), 'temp_images')
        os.makedirs(temp_dir, exist_ok=True)
        
        for slide_number, slide in enumerate(prs.slides):
            for shape in slide.shapes:
                if shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
                    try:
                        image_path = os.path.join(temp_dir, f'slide_{slide_number}.png')
                        with open(image_path, 'wb') as f:
                            f.write(shape.image.blob)
                        images.append(image_path)
                    except Exception as shape_error:
                        print(f"Resim işlenirken hata: {shape_error}")
                        continue
        
        return images
    except Exception as e:
        print(f"PowerPoint'ten resim çıkarılırken hata: {e}")
        return []

def extract_images_from_pdf(file_path):
    """PDF'ten resimleri çıkarır."""
    try:
        images = []
        temp_dir = os.path.join(os.path.dirname(file_path), 'temp_images')
        os.makedirs(temp_dir, exist_ok=True)
        
        # PyMuPDF ile PDF'i aç
        pdf_document = fitz.open(file_path)
        
        for page_num in range(len(pdf_document)):
            page = pdf_document[page_num]
            image_list = page.get_images()
            
            for img_index, img in enumerate(image_list):
                xref = img[0]
                base_image = pdf_document.extract_image(xref)
                image_bytes = base_image["image"]
                
                # Resmi kaydet
                image_path = os.path.join(temp_dir, f'pdf_image_{page_num}_{img_index}.png')
                with open(image_path, 'wb') as image_file:
                    image_file.write(image_bytes)
                
                # Resmi PIL ile aç ve boyutunu kontrol et
                with PILImage.open(image_path) as img:
                    width, height = img.size
                    # Çok küçük resimleri (örn. ikonlar) atla
                    if width > 100 and height > 100:
                        images.append(image_path)
        
        pdf_document.close()
        return images
    except Exception as e:
        print(f"PDF'ten resim çıkarılırken hata: {e}")
        return []

def create_pdf_with_images(text, images, buffer):
    """Metni ve resimleri PDF'e dönüştürür."""
    try:
        doc = SimpleDocTemplate(
            buffer,
            pagesize=letter,
            rightMargin=30,
            leftMargin=30,
            topMargin=30,
            bottomMargin=30
        )

        styles = getSampleStyleSheet()
        normal_style = ParagraphStyle(
            'CustomNormal',
            parent=styles['Normal'],
            fontSize=11,
            leading=14,
            spaceBefore=6,
            spaceAfter=6
        )

        story = []
        paragraphs = text.split("\n\n")

        # Önce metinleri ekle
        for paragraph in paragraphs:
            if paragraph.strip():
                p = Paragraph(paragraph.strip(), normal_style)
                story.append(p)
                story.append(Spacer(1, 10))

        # Sonra resimleri 3x3 grid olarak ekle (her sayfada 9 resim)
        if images:
            # Tekrarlanan resimleri engelle
            unique_images = []
            seen_images = set()
            
            for img_path in images:
                # Resmin hash'ini al (boyut ve ilk birkaç byte'a bakarak)
                try:
                    with open(img_path, 'rb') as f:
                        img_data = f.read()
                        img_hash = hash(img_data[:1000])  # İlk 1000 byte'ı kullan
                        if img_hash not in seen_images:
                            seen_images.add(img_hash)
                            unique_images.append(img_path)
                except:
                    continue

            # Her sayfada 9 resim olacak şekilde grid oluştur
            for i in range(0, len(unique_images), 9):
                page_images = unique_images[i:i+9]
                grid_data = []
                for j in range(0, 9, 3):
                    row = []
                    for img_path in page_images[j:j+3]:
                        try:
                            img = Image(img_path)
                            # Resim boyutunu büyüt
                            img._restrictSize(3*inch, 3*inch)
                            row.append(img)
                        except:
                            row.append('')
                    # Eğer satırda 3'ten az resim varsa boş kutu ekle
                    while len(row) < 3:
                        row.append('')
                    grid_data.append(row)
                
                table = Table(grid_data, colWidths=[3.2*inch]*3)
                table.setStyle(TableStyle([
                    ('ALIGN', (0, 0), (-1, -1), 'CENTER'),
                    ('VALIGN', (0, 0), (-1, -1), 'MIDDLE'),
                    ('LEFTPADDING', (0, 0), (-1, -1), 5),
                    ('RIGHTPADDING', (0, 0), (-1, -1), 5),
                    ('TOPPADDING', (0, 0), (-1, -1), 5),
                    ('BOTTOMPADDING', (0, 0), (-1, -1), 5),
                ]))
                story.append(PageBreak())
                story.append(table)

        doc.build(story)
        return True
    except Exception as e:
        print(f"PDF oluşturulurken hata: {e}")
        return str(e)

def create_word_with_images(text, images, buffer):
    """Metni ve resimleri Word dosyasına dönüştürür."""
    try:
        doc = Document()
        # Metinleri ekle
        paragraphs = text.split("\n\n")
        for paragraph in paragraphs:
            if paragraph.strip():
                doc.add_paragraph(paragraph.strip())
        
        # Sonra resimleri 3x3 grid olarak ekle (her sayfada 9 resim)
        if images:
            from docx.shared import Inches
            # Tekrarlanan resimleri engelle
            unique_images = []
            seen_images = set()
            
            for img_path in images:
                try:
                    with open(img_path, 'rb') as f:
                        img_data = f.read()
                        img_hash = hash(img_data[:1000])
                        if img_hash not in seen_images:
                            seen_images.add(img_hash)
                            unique_images.append(img_path)
                except:
                    continue

            for i in range(0, len(unique_images), 9):
                page_images = unique_images[i:i+9]
                doc.add_page_break()
                table = doc.add_table(rows=3, cols=3)
                table.autofit = True
                for idx, img_path in enumerate(page_images):
                    row = idx // 3
                    col = idx % 3
                    cell = table.cell(row, col)
                    try:
                        run = cell.paragraphs[0].add_run()
                        # Resim boyutunu büyüt
                        run.add_picture(img_path, width=Inches(3))
                    except:
                        cell.text = ''
        
        doc.save(buffer)
        return True
    except Exception as e:
        print(f"Word dosyası oluşturulurken hata: {e}")
        return str(e)

def home(request):
    if request.method == 'POST' and request.FILES['document']:
        uploaded_file = request.FILES['document']
        fs = FileSystemStorage()
        name = fs.save(uploaded_file.name, uploaded_file)
        file_path = fs.path(name)
        
        # Dosya uzantısını kontrol et
        file_ext = os.path.splitext(name)[1].lower()
        
        # Çıktı formatını kontrol et
        output_format = request.POST.get('output_format', 'pdf')
        
        # Resimli dönüştürme isteği mi kontrol et
        with_images = request.POST.get('with_images', '') == 'true'
        
        # Dosya tipine göre metni çıkar
        if file_ext in ['.doc', '.docx']:
            text = extract_text_from_word(file_path)
            images = []
        elif file_ext in ['.ppt', '.pptx']:
            text = extract_text_from_powerpoint(file_path)
            images = extract_images_from_powerpoint(file_path) if with_images else []
        elif file_ext == '.pdf':
            text = extract_text_from_pdf(file_path)
            images = extract_images_from_pdf(file_path) if with_images else []
        else:
            fs.delete(name)
            return HttpResponse("Desteklenmeyen dosya formatı. Lütfen Word (.doc, .docx), PowerPoint (.ppt, .pptx) veya PDF (.pdf) dosyası yükleyin.", status=400)
        
        # Çıktı formatına göre dönüştür
        buffer = io.BytesIO()
        if output_format == 'pdf':
            result = create_pdf_with_images(text, images, buffer)
            content_type = 'application/pdf'
            file_extension = 'pdf'
        else:  # word
            result = create_word_with_images(text, images, buffer)
            content_type = 'application/vnd.openxmlformats-officedocument.wordprocessingml.document'
            file_extension = 'docx'
        
        # Geçici dosyaları temizle
        fs.delete(name)
        if 'temp_dir' in locals():
            for img_path in images:
                try:
                    os.remove(img_path)
                except:
                    pass
            try:
                os.rmdir(os.path.dirname(images[0]))
            except:
                pass
        
        if result is True:
            buffer.seek(0)
            response = FileResponse(buffer, as_attachment=True, filename=f'converted.{file_extension}')
            response['Content-Type'] = content_type
            return response
        else:
            return HttpResponse(f"Dönüştürme sırasında bir hata oluştu: {result}", status=500)
    
    return render(request, 'converter/home.html', {'range': range(9)})